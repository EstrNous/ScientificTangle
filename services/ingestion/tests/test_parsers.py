import subprocess
import zipfile
from io import BytesIO
from pathlib import Path

import fitz
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from app.parsers import ParserRegistry, SourceContent
from shared.contracts import AccessPolicy


def registry() -> ParserRegistry:
    return ParserRegistry(
        libreoffice_binary="soffice",
        conversion_timeout_seconds=5,
        archive_max_entries=10,
        archive_max_uncompressed_bytes=1024 * 1024,
    )


def source(filename: str, content: bytes) -> SourceContent:
    return SourceContent(
        object_key=f"uploads/user/task/{filename}",
        original_filename=filename,
        source_checksum="a" * 64,
        content=content,
    )


def docx_content() -> bytes:
    document = Document()
    document.add_paragraph("Никель показал извлекаемость 82 %.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "parameter"
    table.cell(0, 1).text = "value"
    table.cell(1, 0).text = "recovery"
    table.cell(1, 1).text = "82 %"
    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def test_docx_parser_creates_stable_spans_and_table_rows() -> None:
    parser = registry()
    stored = source("report.docx", docx_content())

    first, first_warnings = parser.normalize(stored, AccessPolicy())
    second, second_warnings = parser.normalize(stored, AccessPolicy())

    assert first_warnings == second_warnings == []
    assert first[0].id == second[0].id
    assert first[0].table_blocks[0].headers == ["parameter", "value"]
    assert first[0].table_blocks[0].rows == [["recovery", "82 %"]]
    assert sum(span.source_type == "table" for span in first[0].source_spans) == 2


def test_pdf_parser_creates_page_spans() -> None:
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Nickel recovery 82 percent")
    content = pdf.tobytes()
    pdf.close()

    documents, warnings = registry().normalize(
        source("report.pdf", content),
        AccessPolicy(),
    )

    assert documents[0].source_spans[0].page == 1
    assert "Nickel recovery" in documents[0].content
    assert not any(warning.startswith("pdf_parse_failed") for warning in warnings)


def test_pptx_parser_creates_slide_and_table_spans() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    text_box.text = "Catholyte circulation"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(2)).table
    table.cell(0, 0).text = "flow"
    table.cell(0, 1).text = "value"
    table.cell(1, 0).text = "optimal"
    table.cell(1, 1).text = "3 m/s"
    buffer = BytesIO()
    presentation.save(buffer)

    documents, warnings = registry().normalize(
        source("slides.pptx", buffer.getvalue()),
        AccessPolicy(),
    )

    assert warnings == []
    assert documents[0].source_spans[0].page == 1
    assert documents[0].table_blocks[0].rows == [["optimal", "3 m/s"]]


def test_doc_parser_uses_libreoffice_conversion(monkeypatch) -> None:
    converted = docx_content()

    def run(args, capture_output, check, timeout):
        output_dir = Path(args[args.index("--outdir") + 1])
        (output_dir / "source.docx").write_bytes(converted)
        return subprocess.CompletedProcess(args, 0)

    monkeypatch.setattr(subprocess, "run", run)

    documents, warnings = registry().normalize(
        source("legacy.doc", b"legacy"),
        AccessPolicy(),
    )

    assert documents[0].metadata["parser"] == "libreoffice+python-docx"
    assert warnings == ["legacy_doc_converted:legacy.doc"]


def test_zip_parser_normalizes_entries_and_rejects_unsafe_paths() -> None:
    safe_buffer = BytesIO()
    with zipfile.ZipFile(safe_buffer, "w") as archive:
        archive.writestr("reports/report.docx", docx_content())
        archive.writestr("notes.txt", "ignored")

    documents, warnings = registry().normalize(
        source("corpus.zip", safe_buffer.getvalue()),
        AccessPolicy(),
    )

    assert len(documents) == 1
    assert documents[0].metadata["archive_entry"] == "reports/report.docx"
    assert warnings == ["unsupported_source_format:notes.txt"]

    unsafe_buffer = BytesIO()
    with zipfile.ZipFile(unsafe_buffer, "w") as archive:
        archive.writestr("../report.docx", docx_content())

    unsafe_documents, unsafe_warnings = registry().normalize(
        source("unsafe.zip", unsafe_buffer.getvalue()),
        AccessPolicy(),
    )

    assert unsafe_documents == []
    assert unsafe_warnings == ["zip_path_rejected:unsafe.zip"]
