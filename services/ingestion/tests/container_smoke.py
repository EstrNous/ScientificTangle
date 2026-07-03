import subprocess
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory

import fitz
from docx import Document
from pptx import Presentation

from app.parsers import ParserRegistry, SourceContent
from shared.contracts import AccessPolicy


def registry() -> ParserRegistry:
    return ParserRegistry(
        libreoffice_binary="soffice",
        conversion_timeout_seconds=30,
        archive_max_entries=10,
        archive_max_uncompressed_bytes=1024 * 1024,
    )


def source(filename: str, content: bytes) -> SourceContent:
    return SourceContent(
        object_key=f"uploads/user/task/{filename}",
        original_filename=filename,
        source_checksum="b" * 64,
        content=content,
    )


def main() -> None:
    parser = registry()
    word = Document()
    word.add_paragraph("Nickel recovery 82 percent")
    docx_buffer = BytesIO()
    word.save(docx_buffer)
    docx_documents, _ = parser.normalize(
        source("report.docx", docx_buffer.getvalue()),
        AccessPolicy(),
    )
    assert docx_documents[0].source_spans

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Catholyte circulation"
    pptx_buffer = BytesIO()
    presentation.save(pptx_buffer)
    pptx_documents, _ = parser.normalize(
        source("slides.pptx", pptx_buffer.getvalue()),
        AccessPolicy(),
    )
    assert pptx_documents[0].source_spans

    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "Mine water injection")
    pdf_documents, _ = parser.normalize(
        source("report.pdf", pdf.tobytes()),
        AccessPolicy(),
    )
    pdf.close()
    assert pdf_documents[0].source_spans

    with TemporaryDirectory(prefix="st-smoke-") as temp_dir:
        docx_path = Path(temp_dir) / "legacy.docx"
        docx_path.write_bytes(docx_buffer.getvalue())
        converted = subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "doc:MS Word 97",
                "--outdir",
                temp_dir,
                str(docx_path),
            ],
            capture_output=True,
            check=False,
            timeout=30,
        )
        doc_path = Path(temp_dir) / "legacy.doc"
        assert converted.returncode == 0 and doc_path.is_file()
        doc_documents, warnings = parser.normalize(
            source("legacy.doc", doc_path.read_bytes()),
            AccessPolicy(),
        )
        assert doc_documents[0].source_spans
        assert warnings == ["legacy_doc_converted:legacy.doc"]


if __name__ == "__main__":
    main()
