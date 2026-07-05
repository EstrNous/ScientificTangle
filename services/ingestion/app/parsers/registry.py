import csv
import hashlib
import json
import stat
import subprocess
import zipfile
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path, PurePosixPath
from tempfile import TemporaryDirectory
from typing import Literal

import fitz
from docx import Document
from pptx import Presentation

from shared.contracts import AccessPolicy, NormalizedDocument, SourceSpan, TableBlock

from ..normalization import enrich_normalized_document


class ParserError(Exception):
    def __init__(self, code: str, message: str) -> None:
        self.code = code
        self.message = message
        super().__init__(message)


@dataclass(frozen=True, slots=True)
class SourceContent:
    object_key: str
    original_filename: str
    source_checksum: str
    content: bytes
    archive_entry: str | None = None

    @property
    def checksum(self) -> str:
        return hashlib.sha256(self.content).hexdigest()


@dataclass(frozen=True, slots=True)
class ParsedBlock:
    page: int
    text: str
    source_type: Literal["text", "table", "figure", "caption"] = "text"
    table_block_id: str | None = None


class ParserRegistry:
    supported_extensions = frozenset({
        ".pdf",
        ".docx",
        ".pptx",
        ".doc",
        ".zip",
        ".txt",
        ".md",
        ".csv",
        ".tsv",
        ".json",
    })

    def __init__(
        self,
        libreoffice_binary: str,
        conversion_timeout_seconds: float,
        archive_max_entries: int,
        archive_max_uncompressed_bytes: int,
    ) -> None:
        self._libreoffice_binary = libreoffice_binary
        self._conversion_timeout_seconds = conversion_timeout_seconds
        self._archive_max_entries = archive_max_entries
        self._archive_max_uncompressed_bytes = archive_max_uncompressed_bytes

    def normalize(
        self,
        source: SourceContent,
        access_policy: AccessPolicy,
    ) -> tuple[list[NormalizedDocument], list[str]]:
        extension = Path(source.original_filename).suffix.lower()
        if extension not in self.supported_extensions:
            return [], [f"unsupported_source_format:{source.original_filename}"]
        if extension == ".zip":
            try:
                return self._normalize_archive(source, access_policy)
            except ParserError as error:
                return [], [f"{error.code}:{source.original_filename}"]
        if extension == ".json":
            return self._parse_json(source, access_policy)
        try:
            document, warnings = self._normalize_document(source, access_policy, extension)
        except ParserError as error:
            return [], [f"{error.code}:{source.original_filename}"]
        except Exception:
            return [], [f"parse_failed:{source.original_filename}"]
        if not document.source_spans and not document.table_blocks:
            warnings.append(f"empty_document:{source.original_filename}")
        return [document], warnings

    def _normalize_archive(
        self,
        source: SourceContent,
        access_policy: AccessPolicy,
    ) -> tuple[list[NormalizedDocument], list[str]]:
        documents: list[NormalizedDocument] = []
        warnings: list[str] = []
        try:
            archive = zipfile.ZipFile(BytesIO(source.content))
        except (zipfile.BadZipFile, OSError) as error:
            raise ParserError("invalid_zip", "ZIP archive is invalid") from error
        with archive:
            entries = [entry for entry in archive.infolist() if not entry.is_dir()]
            if len(entries) > self._archive_max_entries:
                raise ParserError("zip_entry_limit_exceeded", "ZIP archive has too many entries")
            total_size = sum(entry.file_size for entry in entries)
            if total_size > self._archive_max_uncompressed_bytes:
                raise ParserError("zip_size_limit_exceeded", "ZIP archive is too large")
            for entry in entries:
                entry_path = self._safe_archive_path(entry)
                extension = Path(entry_path).suffix.lower()
                if extension == ".zip":
                    warnings.append(f"nested_zip_unsupported:{entry_path}")
                    continue
                if extension not in self.supported_extensions:
                    warnings.append(f"unsupported_source_format:{entry_path}")
                    continue
                try:
                    content = archive.read(entry)
                except (RuntimeError, zipfile.BadZipFile, OSError):
                    warnings.append(f"archive_entry_read_failed:{entry_path}")
                    continue
                nested_source = SourceContent(
                    object_key=source.object_key,
                    original_filename=PurePosixPath(entry_path).name,
                    source_checksum=source.source_checksum,
                    content=content,
                    archive_entry=entry_path,
                )
                try:
                    if extension == ".json":
                        entry_documents, entry_warnings = self._parse_json(
                            nested_source,
                            access_policy,
                        )
                        documents.extend(entry_documents)
                        warnings.extend(entry_warnings)
                        continue
                    document, entry_warnings = self._normalize_document(
                        nested_source,
                        access_policy,
                        extension,
                    )
                except ParserError as error:
                    warnings.append(f"{error.code}:{entry_path}")
                    continue
                except Exception:
                    warnings.append(f"parse_failed:{entry_path}")
                    continue
                documents.append(document)
                warnings.extend(entry_warnings)
                if not document.source_spans and not document.table_blocks:
                    warnings.append(f"empty_document:{entry_path}")
        return documents, warnings

    def _normalize_document(
        self,
        source: SourceContent,
        access_policy: AccessPolicy,
        extension: str,
    ) -> tuple[NormalizedDocument, list[str]]:
        document_id = self._document_id(source)
        if extension == ".pdf":
            blocks, tables, warnings = self._parse_pdf(source.content, document_id)
            parser_name = "pymupdf"
        elif extension == ".docx":
            blocks, tables = self._parse_docx(source.content, document_id)
            warnings = []
            parser_name = "python-docx"
        elif extension == ".pptx":
            blocks, tables = self._parse_pptx(source.content, document_id)
            warnings = []
            parser_name = "python-pptx"
        elif extension in {".txt", ".md"}:
            blocks = self._parse_text(source.content)
            tables = []
            warnings = []
            parser_name = "text"
        elif extension in {".csv", ".tsv"}:
            blocks, tables = self._parse_delimited(source.content, document_id, "\t" if extension == ".tsv" else ",")
            warnings = []
            parser_name = "csv"
        else:
            converted = self._convert_doc(source.content)
            blocks, tables = self._parse_docx(converted, document_id)
            warnings = [f"legacy_doc_converted:{source.original_filename}"]
            parser_name = "libreoffice+python-docx"
        document = self._build_document(
            source,
            document_id,
            parser_name,
            blocks,
            tables,
            access_policy,
        )
        return document, warnings

    def _parse_pdf(
        self,
        content: bytes,
        document_id: str,
    ) -> tuple[list[ParsedBlock], list[TableBlock], list[str]]:
        blocks: list[ParsedBlock] = []
        tables: list[TableBlock] = []
        warnings: list[str] = []
        try:
            pdf = fitz.open(stream=content, filetype="pdf")
        except Exception as error:
            raise ParserError("pdf_parse_failed", "PDF parsing failed") from error
        with pdf:
            for page_index, page in enumerate(pdf, start=1):
                for block in page.get_text("blocks"):
                    text = str(block[4]).strip()
                    if text:
                        blocks.append(ParsedBlock(page=page_index, text=text))
                try:
                    found_tables = page.find_tables().tables
                    for table_index, found_table in enumerate(found_tables):
                        rows = self._clean_rows(found_table.extract())
                        if not rows:
                            continue
                        table = self._table_block(document_id, page_index, table_index, rows)
                        tables.append(table)
                        blocks.extend(self._table_blocks(table))
                except Exception:
                    warnings.append(f"pdf_table_extraction_failed:page_{page_index}")
        return blocks, tables, warnings

    def _parse_text(self, content: bytes) -> list[ParsedBlock]:
        text = self._decode_text(content).strip()
        return [ParsedBlock(page=1, text=text)] if text else []

    def _parse_delimited(
        self,
        content: bytes,
        document_id: str,
        delimiter: str,
    ) -> tuple[list[ParsedBlock], list[TableBlock]]:
        reader = csv.reader(self._decode_text(content).splitlines(), delimiter=delimiter)
        rows = self._clean_rows(list(reader))
        if not rows:
            return [], []
        table = self._table_block(document_id, 1, 0, rows)
        return self._table_blocks(table), [table]

    def _parse_json(
        self,
        source: SourceContent,
        access_policy: AccessPolicy,
    ) -> tuple[list[NormalizedDocument], list[str]]:
        try:
            payload = json.loads(self._decode_text(source.content))
        except json.JSONDecodeError:
            return [], [f"invalid_json:{source.original_filename}"]
        if isinstance(payload, dict) and "normalized_documents" in payload:
            return [
                NormalizedDocument.model_validate(item)
                for item in payload["normalized_documents"]
            ], []
        if isinstance(payload, dict) and "source_spans" in payload:
            return [NormalizedDocument.model_validate(payload)], []
        if isinstance(payload, dict):
            document_id = self._document_id(source)
            content = str(payload.get("content") or "")
            title = str(payload.get("title") or source.original_filename)
            document = self._build_document(
                source,
                document_id,
                "json",
                [ParsedBlock(page=1, text=content)] if content else [],
                [],
                access_policy,
            )
            return [document.model_copy(update={"title": title})], []
        return [], [f"unsupported_json_shape:{source.original_filename}"]

    def _parse_docx(
        self,
        content: bytes,
        document_id: str,
    ) -> tuple[list[ParsedBlock], list[TableBlock]]:
        try:
            document = Document(BytesIO(content))
        except Exception as error:
            raise ParserError("docx_parse_failed", "DOCX parsing failed") from error
        blocks = [
            ParsedBlock(page=1, text=paragraph.text.strip())
            for paragraph in document.paragraphs
            if paragraph.text.strip()
        ]
        tables = []
        for table_index, source_table in enumerate(document.tables):
            rows = self._clean_rows(
                [[cell.text for cell in row.cells] for row in source_table.rows]
            )
            if not rows:
                continue
            table = self._table_block(document_id, 1, table_index, rows)
            tables.append(table)
            blocks.extend(self._table_blocks(table))
        return blocks, tables

    def _parse_pptx(
        self,
        content: bytes,
        document_id: str,
    ) -> tuple[list[ParsedBlock], list[TableBlock]]:
        try:
            presentation = Presentation(BytesIO(content))
        except Exception as error:
            raise ParserError("pptx_parse_failed", "PPTX parsing failed") from error
        blocks: list[ParsedBlock] = []
        tables: list[TableBlock] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            table_index = 0
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    rows = self._clean_rows(
                        [[cell.text for cell in row.cells] for row in shape.table.rows]
                    )
                    if rows:
                        table = self._table_block(
                            document_id,
                            slide_index,
                            table_index,
                            rows,
                        )
                        tables.append(table)
                        blocks.extend(self._table_blocks(table))
                        table_index += 1
                    continue
                if getattr(shape, "has_text_frame", False):
                    text = shape.text.strip()
                    if text:
                        blocks.append(ParsedBlock(page=slide_index, text=text))
        return blocks, tables

    def _convert_doc(self, content: bytes) -> bytes:
        with TemporaryDirectory(prefix="st-doc-") as temp_dir:
            input_path = Path(temp_dir) / "source.doc"
            input_path.write_bytes(content)
            try:
                result = subprocess.run(
                    [
                        self._libreoffice_binary,
                        "--headless",
                        "--convert-to",
                        "docx",
                        "--outdir",
                        temp_dir,
                        str(input_path),
                    ],
                    capture_output=True,
                    check=False,
                    timeout=self._conversion_timeout_seconds,
                )
            except FileNotFoundError as error:
                raise ParserError("doc_converter_unavailable", "LibreOffice is unavailable") from error
            except subprocess.TimeoutExpired as error:
                raise ParserError("doc_conversion_timeout", "DOC conversion timed out") from error
            output_path = Path(temp_dir) / "source.docx"
            if result.returncode != 0 or not output_path.is_file():
                raise ParserError("doc_conversion_failed", "DOC conversion failed")
            return output_path.read_bytes()

    def _build_document(
        self,
        source: SourceContent,
        document_id: str,
        parser_name: str,
        blocks: list[ParsedBlock],
        tables: list[TableBlock],
        access_policy: AccessPolicy,
    ) -> NormalizedDocument:
        content_parts: list[str] = []
        source_spans: list[SourceSpan] = []
        offset = 0
        for block in blocks:
            text = block.text.strip()
            if not text:
                continue
            if content_parts:
                content_parts.append("\n\n")
                offset += 2
            start_offset = offset
            content_parts.append(text)
            offset += len(text)
            source_spans.append(
                SourceSpan(
                    document_id=document_id,
                    page=block.page,
                    start_offset=start_offset,
                    end_offset=offset,
                    text=text,
                    table_block_id=block.table_block_id,
                    source_type=block.source_type,
                )
            )
        metadata = {
            "object_key": source.object_key,
            "original_filename": source.original_filename,
            "sha256": source.checksum,
            "source_sha256": source.source_checksum,
            "parser": parser_name,
        }
        if source.archive_entry is not None:
            metadata["archive_entry"] = source.archive_entry
        document = NormalizedDocument(
            id=document_id,
            source_type=Path(source.original_filename).suffix.lower().lstrip("."),
            title=source.original_filename,
            content="".join(content_parts),
            source_spans=source_spans,
            table_blocks=tables,
            metadata=metadata,
            access_policy=access_policy,
        )
        return enrich_normalized_document(document)

    def _safe_archive_path(self, entry: zipfile.ZipInfo) -> str:
        if entry.flag_bits & 0x1:
            raise ParserError("encrypted_zip_entry", "Encrypted ZIP entries are not supported")
        mode = entry.external_attr >> 16
        if stat.S_ISLNK(mode):
            raise ParserError("zip_symlink_rejected", "ZIP symlinks are not supported")
        path = PurePosixPath(entry.filename.replace("\\", "/"))
        if path.is_absolute() or ".." in path.parts:
            raise ParserError("zip_path_rejected", "ZIP entry path is unsafe")
        return path.as_posix()

    @staticmethod
    def _document_id(source: SourceContent) -> str:
        identity = source.source_checksum
        if source.archive_entry is not None:
            identity = f"{identity}:{source.archive_entry}"
        return hashlib.sha256(identity.encode("utf-8")).hexdigest()[:32]

    @staticmethod
    def _clean_rows(rows: list[list[str | None]]) -> list[list[str]]:
        cleaned = [[str(cell or "").strip() for cell in row] for row in rows]
        return [row for row in cleaned if any(row)]

    @staticmethod
    def _decode_text(content: bytes) -> str:
        for encoding in ("utf-8-sig", "utf-8", "cp1251"):
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content.decode("utf-8", errors="replace")

    @staticmethod
    def _table_block(
        document_id: str,
        page: int,
        table_index: int,
        rows: list[list[str]],
    ) -> TableBlock:
        table_id = hashlib.sha1(
            f"{document_id}:{page}:{table_index}".encode()
        ).hexdigest()[:16]
        return TableBlock(
            id=table_id,
            document_id=document_id,
            page=page,
            headers=rows[0],
            rows=rows[1:],
            metadata={"table_index": table_index},
        )

    @staticmethod
    def _table_blocks(table: TableBlock) -> list[ParsedBlock]:
        rows = [table.headers, *table.rows]
        return [
            ParsedBlock(
                page=table.page,
                text=" | ".join(row),
                source_type="table",
                table_block_id=table.id,
            )
            for row in rows
            if any(row)
        ]
